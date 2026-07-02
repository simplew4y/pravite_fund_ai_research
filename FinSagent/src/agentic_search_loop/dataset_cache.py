"""Parsed document text cache for agentic search.

This cache is deliberately an internal extraction layer, similar to the
PDFTextExtractor cache. It is not a corpus root: tools should continue to show
the original source file path while reading searchable text from JSON cache.
"""

from __future__ import annotations

import hashlib
import json
import re
import sqlite3
from collections import defaultdict
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable, Optional, Tuple


PARSED_TEXT_DIR = "parsed_text"
SUPPORTED_PARSED_EXTENSIONS = {".docx", ".xlsx", ".xlsm", ".pptx", ".ppt", ".odp", ".md", ".markdown"}
PPTX_EXTENSIONS = {".pptx"}
UNSUPPORTED_PRESENTATION_EXTENSIONS = {".ppt", ".odp"}


def _load_json(value: Any, default: Any) -> Any:
    if value in (None, ""):
        return default
    if isinstance(value, (dict, list)):
        return value
    try:
        return json.loads(str(value))
    except Exception:
        return default


def _normalize_title_path(value: Any) -> str:
    payload = _load_json(value, None)
    if isinstance(payload, list):
        return " > ".join(str(item) for item in payload if str(item).strip())
    if isinstance(payload, str):
        return payload
    return str(value or "")


def _markdown_text(value: Any) -> str:
    return ("" if value is None else str(value)).replace("\r\n", "\n").replace("\r", "\n")


def _abs_from_stored(stored_path: str, dataset_root: Path) -> Path:
    path = Path(stored_path or "")
    if path.is_absolute():
        return path.resolve()
    if path.parts and path.parts[0] == "datasets":
        return (dataset_root.parent / path).resolve()
    if path.parts and path.parts[0] == dataset_root.name:
        return (dataset_root.parent / path).resolve()
    return (dataset_root / path).resolve()


def _format_location(location: dict[str, Any]) -> str:
    parts: list[str] = []
    if location.get("page_start") not in (None, ""):
        if location.get("page_end") not in (None, "", location.get("page_start")):
            parts.append(f"pages {location.get('page_start')}-{location.get('page_end')}")
        else:
            parts.append(f"page {location.get('page_start')}")
    if location.get("slide_start") not in (None, ""):
        if location.get("slide_end") not in (None, "", location.get("slide_start")):
            parts.append(f"slides {location.get('slide_start')}-{location.get('slide_end')}")
        else:
            parts.append(f"slide {location.get('slide_start')}")
    if location.get("sheet_name"):
        sheet = str(location.get("sheet_name"))
        cell_range = str(location.get("cell_range") or "")
        parts.append(f"{sheet}!{cell_range}" if cell_range else sheet)
    if location.get("heading_path"):
        parts.append(f"heading={location.get('heading_path')}")
    if location.get("display_text"):
        parts.append(str(location.get("display_text")))

    metadata = _load_json(location.get("metadata_json"), {})
    if isinstance(metadata, dict):
        for key in (
            "paragraph_index",
            "paragraph_start",
            "paragraph_end",
            "table_index",
            "image_index",
            "line_start",
            "line_end",
        ):
            if metadata.get(key) not in (None, ""):
                parts.append(f"{key}={metadata.get(key)}")
    return "; ".join(parts)


def _connect_collection(dataset_root: Path) -> sqlite3.Connection:
    db_path = dataset_root / "meta" / "collection.sqlite3"
    if not db_path.is_file():
        raise FileNotFoundError(f"collection.sqlite3 not found: {db_path}")
    conn = sqlite3.connect(db_path)
    conn.row_factory = sqlite3.Row
    return conn


def _fetch_document_for_source(conn: sqlite3.Connection, dataset_root: Path, source: Path) -> Optional[dict[str, Any]]:
    source = source.resolve()
    rows = conn.execute(
        """
        SELECT *
        FROM documents
        WHERE deleted_at IS NULL
        ORDER BY created_at DESC
        """
    ).fetchall()
    fallback: Optional[dict[str, Any]] = None
    for row in rows:
        doc = dict(row)
        stored = _abs_from_stored(str(doc.get("stored_path") or ""), dataset_root)
        if stored == source:
            return doc
        if stored.name == source.name:
            fallback = doc
    return fallback


def _fetch_chunks_and_locations(
    conn: sqlite3.Connection,
    doc_id: str,
) -> tuple[list[dict[str, Any]], dict[str, list[dict[str, Any]]]]:
    chunks = [
        dict(row)
        for row in conn.execute(
            """
            SELECT *
            FROM chunks
            WHERE doc_id = ?
            ORDER BY chunk_index ASC
            """,
            (doc_id,),
        ).fetchall()
    ]
    locations: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for row in conn.execute(
        """
        SELECT *
        FROM chunk_locations
        WHERE doc_id = ?
        ORDER BY chunk_id ASC, location_index ASC
        """,
        (doc_id,),
    ).fetchall():
        item = dict(row)
        locations[str(item.get("chunk_id"))].append(item)
    return chunks, locations


def _render_chunks(
    doc: dict[str, Any],
    chunks: list[dict[str, Any]],
    locations: dict[str, list[dict[str, Any]]],
) -> str:
    lines: list[str] = [
        f"# Parsed source text: {doc.get('original_filename') or doc.get('title')}",
        "",
        f"original_filename: {doc.get('original_filename')}",
        f"doc_id: {doc.get('doc_id')}",
        f"file_type: {doc.get('file_type')}",
        f"doc_type: {doc.get('doc_type') or ''}",
        f"stored_path: {doc.get('stored_path')}",
        f"status: {doc.get('status')}",
        "",
    ]
    if not chunks:
        lines.extend(["No parsed chunks are registered for this source file.", ""])
        return "\n".join(lines).rstrip() + "\n"

    for chunk in chunks:
        title = _normalize_title_path(chunk.get("title_path")) or f"chunk {chunk.get('chunk_index')}"
        lines.extend(
            [
                f"## Chunk {chunk.get('chunk_index')}: {title}",
                f"chunk_id: {chunk.get('chunk_id')}",
                f"content_type: {chunk.get('content_type')}",
                f"source_ref: {chunk.get('source_ref') or ''}",
                f"summary: {_markdown_text(chunk.get('summary') or '')}",
            ]
        )
        locs = locations.get(str(chunk.get("chunk_id")), [])
        if locs:
            lines.append("locations:")
            for loc in locs:
                summary = _format_location(loc)
                if summary:
                    lines.append(f"- {summary}")
        lines.extend(["", _markdown_text(chunk.get("content") or ""), ""])
    return "\n".join(lines).rstrip() + "\n"


def _iter_pptx_text(path: Path) -> Iterable[tuple[int, str]]:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(str(path))
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and str(text).strip():
                parts.append(str(text).strip())
            if getattr(shape, "has_table", False):
                table = shape.table
                rows: list[str] = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    if len(rows) > 1:
                        separator = "| " + " | ".join("---" for _ in table.columns) + " |"
                        parts.extend([rows[0], separator, *rows[1:]])
                    else:
                        parts.append(rows[0])
        yield slide_idx, "\n\n".join(part for part in parts if part.strip())


def _render_pptx(path: Path, doc: Optional[dict[str, Any]]) -> str:
    lines: list[str] = [
        f"# Parsed source text: {path.name}",
        "",
        f"original_filename: {(doc or {}).get('original_filename') or path.name}",
        f"doc_id: {(doc or {}).get('doc_id') or ''}",
        "file_type: ppt",
        f"stored_path: {(doc or {}).get('stored_path') or str(path)}",
        "extraction: python-pptx",
        "",
    ]
    slide_count = 0
    for slide_idx, text in _iter_pptx_text(path):
        slide_count += 1
        lines.extend([f"## Slide {slide_idx}", f"source_ref: slide {slide_idx}", ""])
        lines.append(text.strip() if text.strip() else "(no extractable slide text)")
        lines.append("")
    if slide_count == 0:
        lines.extend(["No slides were found in this presentation.", ""])
    return "\n".join(lines).rstrip() + "\n"


@dataclass
class ParsedDocumentCache:
    dataset_root: Path

    @property
    def cache_dir(self) -> Path:
        return self.dataset_root / ".agentic_search_cache" / PARSED_TEXT_DIR

    @property
    def collection_db_path(self) -> Path:
        return self.dataset_root / "meta" / "collection.sqlite3"

    def extract(self, path: Path) -> Tuple[list[str], str]:
        path = path.expanduser().resolve()
        ext = path.suffix.lower()
        if ext in UNSUPPORTED_PRESENTATION_EXTENSIONS:
            raise ValueError(f"Unsupported presentation format for agentic search cache: {path.suffix}")
        if ext not in SUPPORTED_PARSED_EXTENSIONS:
            raise ValueError(f"Unsupported parsed document type: {path.suffix}")
        if not path.is_file():
            raise FileNotFoundError(f"Source file not found: {path}")

        stat = path.stat()
        db_mtime_ns = self.collection_db_path.stat().st_mtime_ns if self.collection_db_path.is_file() else 0
        key = hashlib.sha1(
            f"{path}|{stat.st_mtime_ns}|{stat.st_size}|{db_mtime_ns}".encode("utf-8", errors="replace")
        ).hexdigest()
        cache_file = self.cache_dir / f"{key}.json"
        if cache_file.is_file():
            payload = json.loads(cache_file.read_text(encoding="utf-8"))
            text = str(payload.get("text") or "")
            return text.splitlines(), str(payload.get("method") or "dataset_cache")

        text, method, doc = self._extract_uncached(path)
        self.cache_dir.mkdir(parents=True, exist_ok=True)
        payload = {
            "source": str(path),
            "original_filename": (doc or {}).get("original_filename") or path.name,
            "doc_id": (doc or {}).get("doc_id"),
            "file_type": (doc or {}).get("file_type") or ext.lstrip("."),
            "mtime_ns": stat.st_mtime_ns,
            "size": stat.st_size,
            "collection_db_mtime_ns": db_mtime_ns,
            "method": method,
            "text": text,
        }
        cache_file.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
        return text.splitlines(), method

    def _extract_uncached(self, path: Path) -> tuple[str, str, Optional[dict[str, Any]]]:
        doc: Optional[dict[str, Any]] = None
        if self.collection_db_path.is_file():
            with _connect_collection(self.dataset_root) as conn:
                doc = _fetch_document_for_source(conn, self.dataset_root, path)
                if doc:
                    chunks, locations = _fetch_chunks_and_locations(conn, str(doc.get("doc_id")))
                    if chunks:
                        return _render_chunks(doc, chunks, locations), "dataset_cache", doc

        if path.suffix.lower() in PPTX_EXTENSIONS:
            return _render_pptx(path, doc), "pptx_text", doc

        raise ValueError(
            f"No parsed chunks found for source file: {path}. Run the dataset parser/index pipeline first."
        )


def find_dataset_root(start: str | Path) -> Optional[Path]:
    path = Path(start).expanduser().resolve()
    candidates = [path] if path.is_dir() else [path.parent]
    candidates.extend(candidates[0].parents)
    for candidate in candidates:
        if (candidate / "meta" / "collection.sqlite3").is_file():
            return candidate
    return None
