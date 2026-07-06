#!/usr/bin/env python3
"""Step 1: parse PPTX into ordered slide blocks."""

from __future__ import annotations

import argparse
import re
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from common import dump_json, ensure_supported_file, normalize_text, safe_stem

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as exc:  # pragma: no cover - dependency check happens at runtime.
    raise ImportError("python-pptx is required for file2chunk_ppt; install package 'python-pptx'.") from exc


EMU_PER_INCH = 914400
REL_NS = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def _shape_bbox(shape: Any) -> dict[str, Any]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "x0": left,
        "y0": top,
        "x1": left + width,
        "y1": top + height,
        "unit": "emu",
    }


def _shape_name(shape: Any) -> str:
    return str(getattr(shape, "name", "") or "")


def _shape_type_name(shape: Any) -> str:
    try:
        return str(shape.shape_type).split(".")[-1]
    except Exception:
        return str(getattr(shape, "shape_type", "") or "")


def _paragraph_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = normalize_text("".join(run.text for run in paragraph.runs))
        if not text:
            text = normalize_text(paragraph.text)
        if not text:
            continue
        level = int(getattr(paragraph, "level", 0) or 0)
        prefix = "  " * level + "- " if level > 0 else ""
        lines.append(prefix + text)
    return "\n".join(lines).strip()


def _table_to_markdown(shape: Any) -> tuple[str, dict[str, Any]]:
    rows: list[list[str]] = []
    for row in shape.table.rows:
        cells = [normalize_text(cell.text) for cell in row.cells]
        rows.append(cells)
    if not rows:
        return "", {"row_count": 0, "column_count": 0}
    max_cols = max(len(row) for row in rows)
    padded = [row + [""] * (max_cols - len(row)) for row in rows]
    header = padded[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(max_cols)) + " |",
    ]
    for row in padded[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines), {
        "row_count": len(padded),
        "column_count": max_cols,
        "header": header,
    }


def _slide_xml_name(slide_number: int) -> str:
    return f"ppt/slides/slide{slide_number}.xml"


def _part_slide_file_number(slide: Any, fallback: int) -> int:
    try:
        partname = str(slide.part.partname)
    except Exception:
        return fallback
    match = re.search(r"slide(\d+)\.xml$", partname)
    return int(match.group(1)) if match else fallback


def _notes_target_for_slide(zf: zipfile.ZipFile, slide_number: int) -> str | None:
    rels_name = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    if rels_name not in zf.namelist():
        return None
    root = ET.fromstring(zf.read(rels_name))
    for rel in root.findall("rel:Relationship", REL_NS):
        rel_type = rel.attrib.get("Type", "")
        if rel_type.endswith("/notesSlide"):
            target = rel.attrib.get("Target", "")
            if target.startswith("../"):
                return "ppt/" + target[3:]
            if target.startswith("/"):
                return target.lstrip("/")
            return "ppt/slides/" + target
    return None


def _extract_notes(pptx_path: Path) -> dict[int, str]:
    notes: dict[int, str] = {}
    try:
        with zipfile.ZipFile(pptx_path) as zf:
            slide_numbers = sorted(
                int(match.group(1))
                for name in zf.namelist()
                for match in [re.match(r"ppt/slides/slide(\d+)\.xml$", name)]
                if match
            )
            for slide_number in slide_numbers:
                target = _notes_target_for_slide(zf, slide_number)
                if not target or target not in zf.namelist():
                    continue
                root = ET.fromstring(zf.read(target))
                texts = [normalize_text(node.text or "") for node in root.findall(".//a:t", A_NS)]
                text = "\n".join(item for item in texts if item).strip()
                if text:
                    notes[slide_number] = text
    except (zipfile.BadZipFile, ET.ParseError):
        return {}
    return notes


def _is_title_shape(slide: Any, shape: Any) -> bool:
    try:
        return slide.shapes.title is not None and shape == slide.shapes.title
    except Exception:
        return False


def _slide_title(slide: Any, fallback: str) -> str:
    try:
        title_shape = slide.shapes.title
        if title_shape is not None:
            text = _paragraph_text(title_shape)
            if text:
                return normalize_text(text.splitlines()[0])
    except Exception:
        pass
    for shape in slide.shapes:
        text = _paragraph_text(shape)
        if text and len(text) <= 160:
            return normalize_text(text.splitlines()[0])
    return fallback


def parse_ppt_to_blocks(pptx_path: str | Path) -> dict[str, Any]:
    file_path = ensure_supported_file(pptx_path)
    presentation = Presentation(str(file_path))
    notes_by_slide = _extract_notes(file_path)
    stem = safe_stem(file_path)
    blocks: list[dict[str, Any]] = []
    slides: list[dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides):
        slide_number = slide_index + 1
        slide_file_number = _part_slide_file_number(slide, slide_number)
        title = _slide_title(slide, f"Slide {slide_number}")
        heading_path = [title] if title else []
        slide_blocks: list[str] = []
        table_index = 0
        image_index = 0
        for shape_index, shape in enumerate(slide.shapes, start=1):
            bbox = _shape_bbox(shape)
            base_location = {
                "slide_index": slide_index,
                "slide_number": slide_number,
                "slide_file_number": slide_file_number,
                "slide_title": title,
                "shape_index": shape_index,
                "shape_name": _shape_name(shape),
                "shape_type": _shape_type_name(shape),
                "bbox": bbox,
                "heading_path": heading_path,
            }
            if getattr(shape, "has_table", False):
                table_index += 1
                content, analysis = _table_to_markdown(shape)
                if not content.strip():
                    continue
                block_id = f"slide{slide_number:04d}_table{table_index:03d}"
                loc = {
                    **base_location,
                    "table_index": table_index,
                    "display_text": f"slide {slide_number} table {table_index}",
                }
                blocks.append(
                    {
                        "block_id": block_id,
                        "block_type": "table",
                        "content": content,
                        "slide_index": slide_index,
                        "slide_number": slide_number,
                        "slide_title": title,
                        "heading_path": heading_path,
                        "source_location": loc,
                        "table_analysis": analysis,
                    }
                )
                slide_blocks.append(block_id)
                continue

            text = _paragraph_text(shape)
            if text:
                block_type = "title" if _is_title_shape(slide, shape) else "text"
                block_id = f"slide{slide_number:04d}_{block_type}{shape_index:03d}"
                loc = {
                    **base_location,
                    "display_text": f"slide {slide_number} {block_type} shape {shape_index}",
                }
                blocks.append(
                    {
                        "block_id": block_id,
                        "block_type": block_type,
                        "content": text,
                        "slide_index": slide_index,
                        "slide_number": slide_number,
                        "slide_title": title,
                        "heading_path": heading_path,
                        "source_location": loc,
                    }
                )
                slide_blocks.append(block_id)
                continue

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image_index += 1
                block_id = f"slide{slide_number:04d}_image{image_index:03d}"
                loc = {
                    **base_location,
                    "image_index": image_index,
                    "display_text": f"slide {slide_number} image {image_index}",
                }
                blocks.append(
                    {
                        "block_id": block_id,
                        "block_type": "image",
                        "content": "",
                        "slide_index": slide_index,
                        "slide_number": slide_number,
                        "slide_title": title,
                        "heading_path": heading_path,
                        "source_location": loc,
                    }
                )
                slide_blocks.append(block_id)

        notes = notes_by_slide.get(slide_file_number, "")
        if notes:
            block_id = f"slide{slide_number:04d}_notes001"
            loc = {
                "slide_index": slide_index,
                "slide_number": slide_number,
                "slide_file_number": slide_file_number,
                "slide_title": title,
                "shape_index": None,
                "shape_name": "speaker_notes",
                "shape_type": "notes",
                "bbox": None,
                "heading_path": heading_path,
                "notes_index": 1,
                "display_text": f"slide {slide_number} speaker notes",
            }
            blocks.append(
                {
                    "block_id": block_id,
                    "block_type": "notes",
                    "content": notes,
                    "slide_index": slide_index,
                    "slide_number": slide_number,
                    "slide_title": title,
                    "heading_path": heading_path,
                    "source_location": loc,
                }
            )
            slide_blocks.append(block_id)

        slides.append(
            {
                "slide_index": slide_index,
                "slide_number": slide_number,
                "title": title,
                "block_ids": slide_blocks,
                "shape_count": len(slide.shapes),
            }
        )

    return {
        "source_filename": file_path.name,
        "source_path": str(file_path),
        "stem": stem,
        "slide_width_emu": int(presentation.slide_width),
        "slide_height_emu": int(presentation.slide_height),
        "slide_count": len(presentation.slides),
        "slides": slides,
        "blocks": blocks,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Parse PPTX into ordered slide blocks.")
    parser.add_argument("--file", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()
    dump_json(parse_ppt_to_blocks(args.file), args.output)


if __name__ == "__main__":
    main()
